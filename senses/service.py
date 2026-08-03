#!/usr/bin/env python3
"""
Cosmos Senses — usługa percepcji (zmysły Cosmosa).

Każdy zmysł jest OPCJONALNY: usługa startuje z tym, co masz zainstalowane,
a /health mówi Cosmosowi, które zmysły są dostępne.

    słuch    /stt     Whisper (faster-whisper)     pip install faster-whisper
    głos     /tts     Piper                        pip install piper-tts
    wzrok    /detect  YOLO (ultralytics)           pip install ultralytics
    ciało    /pose    MediaPipe (sylwetka/gesty)   pip install mediapipe

Uruchomienie:
    pip install fastapi uvicorn python-multipart
    python senses/service.py            # port 7060

Konfiguracja przez zmienne środowiskowe:
    SENSES_PORT      port usługi (domyślnie 7060)
    WHISPER_MODEL    small | base | medium | large-v3   (domyślnie small)
    WHISPER_DEVICE   cuda | cpu                          (domyślnie auto)
    PIPER_VOICE      ścieżka do głosu .onnx, np. pl_PL-darkman-medium.onnx
    YOLO_MODEL       domyślnie yolo11n.pt (pobiera się automatycznie)
"""

import base64
import io
import os
import tempfile

from fastapi import FastAPI, Request
from fastapi.responses import JSONResponse, Response, StreamingResponse
import uvicorn

app = FastAPI(title="Cosmos Senses")

# ---------------------------------------------------------------------------
# Wykrywanie dostępnych zmysłów (leniwa inicjalizacja modeli)
# ---------------------------------------------------------------------------

CAPS = {"whisper": False, "piper": False, "yolo": False, "mediapipe": False,
        "embed": False, "upscale": False, "kinect": False}

# Kinect 360 przez oficjalne SDK (tylko Windows). Obraz z niego nie jest widoczny
# dla przeglądarki ani OpenCV — nie jest kamerą UVC — więc klatki muszą iść
# do Cosmosa tędy, przez HTTP.
try:
    import kinect_win
    CAPS["kinect"] = kinect_win.sensor_count() > 0
except Exception:
    pass

try:
    import faster_whisper  # noqa: F401
    CAPS["whisper"] = True
except ImportError:
    pass

try:
    import piper  # noqa: F401
    CAPS["piper"] = bool(os.environ.get("PIPER_VOICE"))
except ImportError:
    pass

try:
    import ultralytics  # noqa: F401
    CAPS["yolo"] = True
except ImportError:
    pass

try:
    import mediapipe  # noqa: F401
    CAPS["mediapipe"] = True
except ImportError:
    pass

try:
    import sentence_transformers  # noqa: F401
    CAPS["embed"] = True
except ImportError:
    pass

try:
    import realesrgan  # noqa: F401
    CAPS["upscale"] = True
except ImportError:
    pass

_whisper_model = None
_piper_voice = None
_yolo_model = None
_embed_model = None


def get_whisper():
    global _whisper_model
    if _whisper_model is None:
        from faster_whisper import WhisperModel
        name = os.environ.get("WHISPER_MODEL", "small")
        device = os.environ.get("WHISPER_DEVICE", "auto")
        compute = "float16" if device != "cpu" else "int8"
        try:
            _whisper_model = WhisperModel(name, device=device, compute_type=compute)
        except Exception:
            _whisper_model = WhisperModel(name, device="cpu", compute_type="int8")
    return _whisper_model


def get_piper():
    global _piper_voice
    if _piper_voice is None:
        from piper import PiperVoice
        _piper_voice = PiperVoice.load(os.environ["PIPER_VOICE"])
    return _piper_voice


def get_yolo():
    global _yolo_model
    if _yolo_model is None:
        from ultralytics import YOLO
        _yolo_model = YOLO(os.environ.get("YOLO_MODEL", "yolo11n.pt"))
    return _yolo_model


def get_embedder():
    global _embed_model
    if _embed_model is None:
        from sentence_transformers import SentenceTransformer
        # bge-m3: bardzo dobre wielojęzyczne embeddingi (ok. 2 GB).
        # Lżejsza alternatywa: paraphrase-multilingual-MiniLM-L12-v2 (~120 MB).
        name = os.environ.get("EMBED_MODEL", "BAAI/bge-m3")
        _embed_model = SentenceTransformer(name)
    return _embed_model


def decode_image(payload: dict):
    """dataURL/base64 -> obraz OpenCV (numpy BGR)."""
    import cv2
    import numpy as np
    data = payload.get("image", "")
    if "," in data:  # data:image/jpeg;base64,....
        data = data.split(",", 1)[1]
    raw = base64.b64decode(data)
    arr = np.frombuffer(raw, dtype=np.uint8)
    return cv2.imdecode(arr, cv2.IMREAD_COLOR)


# ---------------------------------------------------------------------------
# Endpointy
# ---------------------------------------------------------------------------

@app.get("/health")
def health():
    return CAPS


@app.post("/stt")
async def stt(request: Request):
    """Audio (webm/ogg/wav/mp3) w body -> {"text": "..."}"""
    if not CAPS["whisper"]:
        return JSONResponse({"error": "Whisper niezainstalowany (pip install faster-whisper)."}, status_code=501)
    audio = await request.body()
    suffix = ".webm"
    ctype = request.headers.get("content-type", "")
    for ext in ("wav", "ogg", "mp3", "mp4", "m4a"):
        if ext in ctype:
            suffix = "." + ext
            break
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as f:
        f.write(audio)
        tmp = f.name
    try:
        segments, info = get_whisper().transcribe(tmp, language=os.environ.get("WHISPER_LANG") or None, vad_filter=True)
        text = " ".join(s.text.strip() for s in segments).strip()
        return {"text": text, "language": info.language}
    finally:
        os.unlink(tmp)


@app.post("/tts")
async def tts(request: Request):
    """{"text": "..."} -> audio/wav"""
    if not CAPS["piper"]:
        return JSONResponse(
            {"error": "Piper niedostępny (pip install piper-tts + ustaw PIPER_VOICE na plik głosu .onnx)."},
            status_code=501,
        )
    payload = await request.json()
    text = (payload.get("text") or "").strip()
    if not text:
        return JSONResponse({"error": "Puste pole text."}, status_code=400)
    import wave
    buf = io.BytesIO()
    with wave.open(buf, "wb") as wav_file:
        get_piper().synthesize(text, wav_file)
    return Response(content=buf.getvalue(), media_type="audio/wav")


@app.post("/detect")
async def detect(request: Request):
    """{"image": dataURL} -> {"objects": [{label, conf, box}], "summary": "..."}"""
    if not CAPS["yolo"]:
        return JSONResponse({"error": "YOLO niezainstalowany (pip install ultralytics)."}, status_code=501)
    payload = await request.json()
    img = decode_image(payload)
    if img is None:
        return JSONResponse({"error": "Nie udało się zdekodować obrazu."}, status_code=400)
    results = get_yolo()(img, verbose=False)[0]
    objects = []
    for b in results.boxes:
        objects.append({
            "label": results.names[int(b.cls)],
            "conf": round(float(b.conf), 3),
            "box": [round(float(v)) for v in b.xyxy[0].tolist()],
        })
    labels = {}
    for o in objects:
        labels[o["label"]] = labels.get(o["label"], 0) + 1
    summary = ", ".join(f"{v}× {k}" for k, v in sorted(labels.items(), key=lambda x: -x[1])) or "brak wykrytych obiektów"
    return {"objects": objects, "summary": summary}


@app.post("/pose")
async def pose(request: Request):
    """{"image": dataURL} -> {"present": bool, "summary": "..."}"""
    if not CAPS["mediapipe"]:
        return JSONResponse({"error": "MediaPipe niezainstalowany (pip install mediapipe)."}, status_code=501)
    import cv2
    import mediapipe as mp
    payload = await request.json()
    img = decode_image(payload)
    if img is None:
        return JSONResponse({"error": "Nie udało się zdekodować obrazu."}, status_code=400)
    with mp.solutions.pose.Pose(static_image_mode=True) as pose_model:
        res = pose_model.process(cv2.cvtColor(img, cv2.COLOR_BGR2RGB))
    if not res.pose_landmarks:
        return {"present": False, "summary": "nie widać sylwetki"}
    lm = res.pose_landmarks.landmark
    nose_y = lm[0].y
    hip_y = (lm[23].y + lm[24].y) / 2
    posture = "stoi" if (hip_y - nose_y) > 0.45 else "siedzi lub jest blisko kamery"
    return {"present": True, "summary": f"widoczna sylwetka, osoba prawdopodobnie {posture}"}


@app.post("/extract")
async def extract(request: Request):
    """{"name": "plik.xlsx", "data": base64} -> {"text": "..."}
    Wyciąga tekst z PDF/DOCX/XLSX/PPTX na potrzeby bazy wiedzy Cosmosa."""
    payload = await request.json()
    name = str(payload.get("name", ""))
    try:
        data = base64.b64decode(payload.get("data", ""))
    except Exception:
        return JSONResponse({"error": "Nieprawidłowe dane base64."}, status_code=400)
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""

    try:
        if ext == "pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)
        elif ext == "docx":
            import docx
            document = docx.Document(io.BytesIO(data))
            parts = [p.text for p in document.paragraphs]
            for table in document.tables:
                for row in table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
            text = "\n".join(parts)
        elif ext in ("xlsx", "xlsm"):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"## Arkusz: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        lines.append(" | ".join(cells))
                    if len(lines) > 8000:
                        break
            text = "\n".join(lines)
        elif ext == "pptx":
            from pptx import Presentation
            pres = Presentation(io.BytesIO(data))
            parts = []
            for i, slide in enumerate(pres.slides, 1):
                parts.append(f"## Slajd {i}")
                for shape in slide.shapes:
                    if getattr(shape, "text", ""):
                        parts.append(shape.text)
            text = "\n".join(parts)
        else:
            text = data.decode("utf-8", errors="ignore")
        return {"text": text[:200000]}
    except ImportError as e:
        return JSONResponse(
            {"error": f"Brak biblioteki do formatu .{ext} — pip install {e.name}"},
            status_code=501,
        )
    except Exception as e:
        return JSONResponse({"error": f"Błąd ekstrakcji: {e}"}, status_code=400)


@app.post("/upscale")
async def upscale(request: Request):
    """{"image": dataURL, "scale": 4} -> {"image": dataURL} — powiększanie Real-ESRGAN.
    Opcjonalne: pip install realesrgan basicsr  (wymaga GPU dla sensownej szybkości)."""
    try:
        from realesrgan import RealESRGANer  # noqa: F401
        from basicsr.archs.rrdbnet_arch import RRDBNet
    except ImportError:
        return JSONResponse(
            {"error": "Upscale niedostępny — pip install realesrgan basicsr (senses/README.md)."},
            status_code=501,
        )
    import cv2
    import numpy as np
    payload = await request.json()
    img = decode_image(payload)
    if img is None:
        return JSONResponse({"error": "Nieprawidłowy obraz."}, status_code=400)
    scale = int(payload.get("scale", 4))
    model = RRDBNet(num_in_ch=3, num_out_ch=3, num_feat=64, num_block=23, num_grow_ch=32, scale=4)
    up = RealESRGANer(scale=4, model_path=os.environ.get("REALESRGAN_MODEL", "RealESRGAN_x4plus.pth"), model=model)
    out, _ = up.enhance(img, outscale=scale)
    ok, buf = cv2.imencode(".png", out)
    b64 = base64.b64encode(buf.tobytes()).decode()
    return {"image": f"data:image/png;base64,{b64}"}


# ---------------------------------------------------------------------------
# Kinect 360 — klatki po HTTP
# ---------------------------------------------------------------------------
#
# Przeglądarka nie widzi Kinecta (nie jest kamerą UVC), więc podgląd w Cosmosie
# nie może użyć getUserMedia. Zamiast tego serwujemy pojedyncze klatki JPEG,
# a interfejs odświeża je jak zwykły obrazek.

_kinect = None
_kinect_err = ""


def get_kinect():
    """Jedna instancja czujnika na cały proces — Kinect nie znosi dwóch naraz."""
    global _kinect, _kinect_err
    if _kinect is None:
        import kinect_win
        _kinect = kinect_win.Kinect(color=True, depth=True, skeleton=True)
        _kinect.open()
        _kinect_err = ""
    return _kinect


def _to_jpeg(img, quality: int = 80) -> bytes:
    import cv2
    ok, buf = cv2.imencode(".jpg", img, [int(cv2.IMWRITE_JPEG_QUALITY), quality])
    if not ok:
        raise RuntimeError("Nie udało się zakodować JPEG.")
    return buf.tobytes()


@app.get("/kinect/status")
async def kinect_status():
    """Czy czujnik jest dostępny i co potrafi."""
    try:
        import kinect_win
    except Exception as e:
        return {"available": False, "reason": f"brak modułu kinect_win: {e}"}
    n = kinect_win.sensor_count()
    if n <= 0:
        return {"available": False, "reason": "nie widzę czujnika (zasilacz? sterownik SDK 1.8?)"}
    return {"available": True, "sensors": n, "streams": ["color", "depth"], "error": _kinect_err}


def _render(k, stream: str):
    """Klatka gotowa do zakodowania: obraz BGR albo pokolorowana mapa głębi.

    Głębia to milimetry — dla oka mapujemy zasięg 0,5–4 m na paletę. Zera, czyli
    „nie wiem" (cień podczerwieni, szkło, poza zasięgiem), zostają czarne, żeby
    nie udawały pomiaru, którego nie ma.
    """
    if stream == "depth":
        import cv2
        import numpy as np
        frame = k.depth_frame()
        if frame is None:
            return None
        vis = np.clip((frame.astype(np.float32) - 500) / (4000 - 500), 0, 1)
        vis = (vis * 255).astype(np.uint8)
        vis = cv2.applyColorMap(vis, cv2.COLORMAP_TURBO)
        vis[frame == 0] = 0
        return vis
    return k.color_frame()


@app.get("/kinect/stream")
async def kinect_stream(stream: str = "color", fps: int = 15, quality: int = 70):
    """Ciągły strumień MJPEG.

    Pojedyncze klatki przez /kinect/frame znaczą jedno żądanie HTTP na klatkę.
    Przy drodze telefon → VPS → Tailscale → komputer domowy sam obieg zjada
    ćwierć sekundy, więc podgląd klatkuje niezależnie od tego, jak szybki jest
    czujnik. Tutaj połączenie jest jedno, a klatki lecą w nim jedna za drugą —
    przeglądarka odtwarza to natywnie w zwykłym <img>.
    """
    try:
        k = get_kinect()
    except Exception as e:
        return JSONResponse({"error": f"Kinect niedostępny: {e}"}, status_code=503)
    try:
        import cv2  # noqa: F401
    except ImportError:
        return JSONResponse({"error": "Strumień wymaga: pip install opencv-python"},
                            status_code=501)

    delay = 1.0 / max(1, min(30, fps))
    q = max(20, min(95, quality))

    def frames():
        import time as _t
        while True:
            start = _t.time()
            try:
                img = _render(k, stream)
            except Exception:
                break                       # czujnik zniknął — zamknij strumień
            if img is not None:
                jpg = _to_jpeg(img, q)
                yield (b"--frame\r\nContent-Type: image/jpeg\r\n"
                       b"Content-Length: " + str(len(jpg)).encode() + b"\r\n\r\n"
                       + jpg + b"\r\n")
            left = delay - (_t.time() - start)
            if left > 0:
                _t.sleep(left)

    return StreamingResponse(frames(),
                             media_type="multipart/x-mixed-replace; boundary=frame",
                             headers={"Cache-Control": "no-store"})


@app.get("/kinect/frame")
async def kinect_frame(stream: str = "color"):
    """Pojedyncza klatka jako JPEG. `stream` = color albo depth."""
    global _kinect, _kinect_err
    try:
        k = get_kinect()
    except Exception as e:
        _kinect_err = str(e)
        return JSONResponse({"error": f"Kinect niedostępny: {e}"}, status_code=503)

    try:
        img = _render(k, stream)
        if img is None:
            return JSONResponse({"error": "Brak klatki z Kinecta."}, status_code=503)
        return Response(content=_to_jpeg(img), media_type="image/jpeg",
                        headers={"Cache-Control": "no-store"})
    except ImportError:
        return JSONResponse({"error": "Podgląd wymaga: pip install opencv-python"},
                            status_code=501)
    except Exception as e:
        # Czujnik mógł zostać odłączony — następne żądanie spróbuje otworzyć od nowa.
        _kinect_err = str(e)
        try:
            k.close()
        except Exception:
            pass
        _kinect = None
        return JSONResponse({"error": f"Błąd odczytu z Kinecta: {e}"}, status_code=503)


@app.post("/embed")
async def embed(request: Request):
    """{"texts": ["...", ...]} -> {"vectors": [[...], ...]} (pamięć długotrwała)"""
    if not CAPS["embed"]:
        return JSONResponse({"error": "Embeddingi niedostępne (pip install sentence-transformers)."}, status_code=501)
    payload = await request.json()
    texts = payload.get("texts") or []
    if not isinstance(texts, list) or not texts:
        return JSONResponse({"error": "Pole texts (lista) jest wymagane."}, status_code=400)
    vectors = get_embedder().encode([str(t)[:4000] for t in texts], normalize_embeddings=True)
    return {"vectors": [v.tolist() for v in vectors]}


if __name__ == "__main__":
    port = int(os.environ.get("SENSES_PORT", 7060))
    active = ", ".join(k for k, v in CAPS.items() if v) or "brak (zainstaluj zależności)"
    print(f"\n  ✦ Cosmos Senses — port {port}\n  → aktywne zmysły: {active}\n")
    uvicorn.run(app, host="0.0.0.0", port=port, log_level="warning")
